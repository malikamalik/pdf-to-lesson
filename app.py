import os
import io
import base64
import uuid
import json
import pdfplumber
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from flask import Flask, render_template, request, jsonify, send_from_directory
from werkzeug.utils import secure_filename
from anthropic import Anthropic
from dotenv import load_dotenv

load_dotenv()

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 100 * 1024 * 1024  # 100MB max for PPTX

# Vercel has read-only filesystem — use /tmp there
IS_VERCEL = os.environ.get("VERCEL") == "1"
if IS_VERCEL:
    app.config["UPLOAD_FOLDER"] = "/tmp"
else:
    app.config["UPLOAD_FOLDER"] = os.path.join(os.path.dirname(__file__), "output")
    os.makedirs(app.config["UPLOAD_FOLDER"], exist_ok=True)

ALLOWED_EXTENSIONS = {"pdf", "pptx", "ppt"}


def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


def get_file_ext(filename):
    return filename.rsplit(".", 1)[1].lower() if "." in filename else ""


# ==================== PPTX EXTRACTION ====================

def extract_pptx_text(filepath):
    """Extract text from every slide in a PPTX, preserving slide structure."""
    prs = Presentation(filepath)
    full_text = []
    for i, slide in enumerate(prs.slides):
        slide_texts = []
        # Get slide title
        if slide.shapes.title:
            slide_texts.append(f"Title: {slide.shapes.title.text}")
        # Get all text from shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
            # Get text from tables
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        slide_texts.append(f"[Table row] {row_text}")
        # Get speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_texts.append(f"[Speaker Notes] {notes}")
        if slide_texts:
            full_text.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_texts))
    return "\n\n".join(full_text)


def extract_pptx_images(filepath, max_images=50):
    """Extract high-quality embedded images from PPTX with rich contextual metadata.

    For each image, captures: slide title, all text on the slide, shape name,
    image dimensions, and classifies the image type (chart, diagram, photo, etc.)
    so that Claude can place images intelligently in the lesson.
    """
    import hashlib
    from pptx.util import Emu
    prs = Presentation(filepath)
    raw_images = []
    seen_hashes = set()

    # First pass: collect all image blobs with rich context
    all_blobs = []
    for i, slide in enumerate(prs.slides):
        # Get slide title
        slide_title = ""
        if slide.shapes.title:
            slide_title = slide.shapes.title.text.strip()

        # Get ALL text on this slide for context
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                    if row_text:
                        slide_texts.append(f"[Table] {row_text}")

        # Fallback title from first meaningful text
        if not slide_title and slide_texts:
            for t in slide_texts:
                if len(t) > 3:
                    slide_title = t[:80]
                    break

        slide_context = " | ".join(slide_texts[:10])  # first 10 text items for context
        if len(slide_context) > 400:
            slide_context = slide_context[:400] + "..."

        def process_shape(shape, slide_idx, title, context):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    blob = image.blob
                    # Get image dimensions in the slide (EMUs -> inches)
                    w_inches = shape.width / 914400 if shape.width else 0
                    h_inches = shape.height / 914400 if shape.height else 0
                    # Try to get alt text from XML
                    alt_text = ""
                    try:
                        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main",
                                 "p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
                        cNvPr = shape._element.find('.//p:nvPicPr/p:cNvPr', nsmap)
                        if cNvPr is None:
                            cNvPr = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetml}cNvPr')
                        if cNvPr is not None:
                            alt_text = cNvPr.get("descr", "") or ""
                    except Exception:
                        pass

                    all_blobs.append({
                        "blob": blob,
                        "hash": hashlib.md5(blob).hexdigest(),
                        "size": len(blob),
                        "content_type": image.content_type,
                        "slide": slide_idx + 1,
                        "slide_title": title,
                        "slide_context": context,
                        "shape_name": shape.name or "",
                        "alt_text": alt_text,
                        "width_inches": round(w_inches, 1),
                        "height_inches": round(h_inches, 1),
                    })
                except Exception:
                    pass
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    for sub in shape.shapes:
                        if sub.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            try:
                                image = sub.image
                                blob = image.blob
                                w_inches = sub.width / 914400 if sub.width else 0
                                h_inches = sub.height / 914400 if sub.height else 0
                                all_blobs.append({
                                    "blob": blob,
                                    "hash": hashlib.md5(blob).hexdigest(),
                                    "size": len(blob),
                                    "content_type": image.content_type,
                                    "slide": slide_idx + 1,
                                    "slide_title": title,
                                    "slide_context": context,
                                    "shape_name": sub.name or "",
                                    "alt_text": "",
                                    "width_inches": round(w_inches, 1),
                                    "height_inches": round(h_inches, 1),
                                })
                            except Exception:
                                pass
                except Exception:
                    pass

        for shape in slide.shapes:
            process_shape(shape, i, slide_title, slide_context)

    # Count how often each hash appears (to identify repeated decorative elements)
    hash_counts = {}
    for b in all_blobs:
        hash_counts[b["hash"]] = hash_counts.get(b["hash"], 0) + 1

    # Second pass: filter and keep only meaningful images
    MIN_SIZE = 15_000       # Skip images smaller than 15KB (icons, bullets)
    MAX_REPEATS = 3         # Skip images that appear more than 3 times (decorative)

    for b in all_blobs:
        if len(raw_images) >= max_images:
            break
        if b["size"] < MIN_SIZE:
            continue
        if hash_counts[b["hash"]] > MAX_REPEATS:
            continue
        if b["hash"] in seen_hashes:
            continue
        # Skip decorative gradients/rasterized backgrounds
        alt_lower = (b.get("alt_text") or "").lower()
        if any(skip in alt_lower for skip in ["rasterized", "gradient", "background", "/tmp/"]):
            continue
        # Skip very wide/thin banners (likely decorative bars)
        w_in = b["width_inches"]
        h_in = b["height_inches"]
        if h_in > 0 and w_in / h_in > 5:
            continue
        seen_hashes.add(b["hash"])

        b64 = base64.b64encode(b["blob"]).decode("utf-8")
        data_uri = f"data:{b['content_type']};base64,{b64}"

        # Classify image type from shape name and dimensions
        shape_name_lower = b["shape_name"].lower()
        w, h = b["width_inches"], b["height_inches"]
        area = w * h

        if "chart" in shape_name_lower:
            img_type = "chart/graph"
        elif "diagram" in shape_name_lower:
            img_type = "diagram"
        elif "screenshot" in shape_name_lower:
            img_type = "screenshot"
        elif "logo" in shape_name_lower:
            img_type = "logo"
        elif "photo" in shape_name_lower or "picture" in shape_name_lower:
            img_type = "photo"
        elif area > 20:
            img_type = "large illustration/diagram"
        elif area > 8:
            img_type = "illustration"
        elif w > 1.5 * h + 1:
            img_type = "banner/wide graphic"
        elif h > 1.5 * w + 1:
            img_type = "tall graphic/infographic"
        else:
            img_type = "image"

        # Build a rich description
        desc_parts = []
        desc_parts.append(f"From slide {b['slide']}")
        if b["slide_title"]:
            desc_parts.append(f"titled \"{b['slide_title']}\"")
        desc_parts.append(f"[{img_type}, {w}\"x{h}\"]")
        if b["alt_text"]:
            desc_parts.append(f"Alt text: \"{b['alt_text']}\"")
        if b["shape_name"] and b["shape_name"] not in ("Picture", "Image"):
            desc_parts.append(f"Shape: \"{b['shape_name']}\"")
        if b["slide_context"]:
            ctx = b["slide_context"][:150]
            desc_parts.append(f"Context: {ctx}")

        desc = " — ".join(desc_parts)

        raw_images.append({
            "page": b["slide"],
            "data_uri": data_uri,
            "desc": desc,
            "source": "pptx",
            "size": b["size"],
            "slide_title": b["slide_title"],
            "slide_context": b["slide_context"],
            "img_type": img_type,
        })

    # Sort by slide number, then by size (larger = more important) within same slide
    raw_images.sort(key=lambda x: (x["page"], -x["size"]))

    print(f"  PPTX image extraction: {len(all_blobs)} total → {len(raw_images)} kept (filtered {len(all_blobs) - len(raw_images)} duplicates/icons)")
    for idx, img in enumerate(raw_images):
        print(f"    [{idx}] {img['desc'][:120]}")
    return raw_images


# ==================== PDF EXTRACTION (kept for backwards compat) ====================

def extract_pdf_text(filepath):
    """Extract text from PDF using pdfplumber."""
    full_text = []
    with pdfplumber.open(filepath) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text:
                full_text.append(f"--- Page {i + 1} ---\n{text}")
    return "\n\n".join(full_text)


def extract_pdf_images(filepath, max_images=30):
    """Extract images from PDF pages."""
    images = []
    try:
        with pdfplumber.open(filepath) as pdf:
            for i, page in enumerate(pdf.pages):
                if len(images) >= max_images:
                    break
                pil_img = page.to_image(resolution=200).original
                buf = io.BytesIO()
                pil_img.save(buf, format="JPEG", quality=85)
                b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
                images.append({
                    "page": i + 1,
                    "data_uri": f"data:image/jpeg;base64,{b64}",
                    "desc": f"Page {i + 1} screenshot"
                })
    except Exception as e:
        print(f"PDF image extraction warning: {e}")
    return images


def process_uploaded_images(files):
    """Process manually uploaded images into base64 data URIs."""
    images = []
    for f in files:
        if not f or not f.filename:
            continue
        ext = get_file_ext(f.filename)
        if ext not in ("png", "jpg", "jpeg", "gif", "webp", "svg"):
            continue
        blob = f.read()
        content_type = f.content_type or f"image/{ext}"
        b64 = base64.b64encode(blob).decode("utf-8")
        images.append({
            "page": 0,
            "data_uri": f"data:{content_type};base64,{b64}",
            "desc": f"Uploaded: {f.filename}",
            "source": "upload"
        })
    return images


# Ask Claude ONLY for structured JSON slides — NOT the whole HTML
SLIDES_SYSTEM_PROMPT = r"""You are an expert instructional designer. Convert PDF content into a structured JSON array of interactive lesson slides.

OUTPUT FORMAT: Return ONLY a valid JSON array. No markdown fences, no explanation, no text before or after the JSON.

Each slide is an object with these fields:
- "cat": category string (e.g. "Introduction", "Core Concepts", "Knowledge Check", "Interactive Activity", "Milestone", "Common Mistakes", "Review", "Reference", "Completion")
- "t": slide title
- "s": subtitle
- "narration": natural TTS narration text (1-3 sentences, conversational)
- "type": one of "content", "quiz", "matching", "prompt_builder", "ordering", "milestone", "completion"
- "body": the slide body (structure depends on type, see below)

### SLIDE TYPE SCHEMAS:

**"content"** — informational slides:
```json
{
  "type": "content",
  "body": {
    "blocks": [
      {"kind": "text", "html": "<strong>Bold</strong> and normal text..."},
      {"kind": "bullets", "items": ["First point", "Second point"]},
      {"kind": "icons", "items": [{"icon": "emoji", "label": "Label", "desc": "Short desc"}]},
      {"kind": "steps", "items": [{"text": "Step description with <strong>bold</strong>"}]},
      {"kind": "table", "headers": ["Col1","Col2"], "rows": [["cell","cell"]]},
      {"kind": "tip", "label": "Pro Tip", "text": "Tip content..."},
      {"kind": "code", "text": "code or example text"},
      {"kind": "compare", "good_label": "Do This", "good": "Good example", "bad_label": "Not This", "bad": "Bad example"},
      {"kind": "image", "image_idx": 0, "alt": "Description of the image"}
    ]
  }
}
```

**"quiz"** — multiple choice (MUST have exactly 4 options, 1 correct):
```json
{
  "type": "quiz",
  "body": {
    "question": "Question text?",
    "options": ["Option A", "Option B", "Option C", "Option D"],
    "correct": 1,
    "explanations": {"correct": "Why this is right...", "wrong": "The correct answer is B because..."}
  }
}
```

**"matching"** — match 5 pairs:
```json
{
  "type": "matching",
  "body": {
    "pairs": [
      {"left": "Term 1", "right": "Definition 1"},
      {"left": "Term 2", "right": "Definition 2"},
      {"left": "Term 3", "right": "Definition 3"},
      {"left": "Term 4", "right": "Definition 4"},
      {"left": "Term 5", "right": "Definition 5"}
    ]
  }
}
```

**"prompt_builder"** — user builds a response using chips:
```json
{
  "type": "prompt_builder",
  "body": {
    "instructions": "Build a query/response using the chips below...",
    "chips": ["chip 1", "chip 2", "chip 3", "chip 4", "chip 5", "chip 6"],
    "placeholder": "Compose your answer here..."
  }
}
```

**"ordering"** — put steps in correct order:
```json
{
  "type": "ordering",
  "body": {
    "instructions": "Arrange these steps in the correct order:",
    "correct_order": ["First step", "Second step", "Third step", "Fourth step", "Fifth step"]
  }
}
```

**"milestone"** — section celebration:
```json
{
  "type": "milestone",
  "body": {
    "emoji": "🎯",
    "message": "You've mastered the basics!",
    "lessons_done": 4
  }
}
```

**"completion"** — final slide:
```json
{
  "type": "completion",
  "body": {
    "takeaways": ["Key takeaway 1", "Key takeaway 2", "Key takeaway 3", "Key takeaway 4"],
    "cta": "Now go apply what you learned!"
  }
}
```

### CONTENT RULES:
1. Create as many slides as needed to cover ALL content — do NOT skip, summarize, or cut any material from the PDF. Every concept, every detail, every example must be included.
2. KEEP EACH SLIDE CONCISE — aim for 3-5 content blocks max per slide so it fits on a mobile screen without excessive scrolling. If a topic is long, split it across multiple slides rather than cramming it all into one.
3. Never have more than 3 content slides without an interactive element (quiz/matching/ordering/prompt_builder)
4. Include at minimum: 4 quizzes, 1 matching game, 1 prompt builder, 1 ordering exercise
5. Add milestones between major sections
6. Include a "Common Mistakes" content slide
7. Include a "Review" or cheat-sheet content slide before completion
8. End with a completion slide
9. Add compare blocks (do/don't) where relevant
10. Add tip blocks throughout
11. Make quiz questions challenging but fair
12. CRITICAL: Preserve ALL content from the PDF. Do not truncate, summarize, or omit any information. If the PDF has 20 topics, make slides for all 20. Every key point, example, statistic, and detail must appear in the lesson.
13. IMAGE PLACEMENT — this is CRITICAL:
    - Each image has a description telling you which ORIGINAL slide it came from and what content was on that slide.
    - You MUST place each image in the lesson slide that covers the SAME TOPIC as the original slide the image came from.
    - Match images by their slide context/title to the lesson content. For example, if image_idx 3 is "From slide 8 titled 'Revenue Model'" then place it in your lesson slide about revenue.
    - Use {"kind": "image", "image_idx": N, "alt": "DESCRIPTIVE LABEL"} blocks.
    - The "alt" MUST be a specific, descriptive label like "Revenue growth chart showing 3x increase" or "Product architecture diagram" — NOT generic labels like "Image from slide 5".
    - Place images AFTER the introductory text for that topic, so the reader understands the context before seeing the image.
    - Every image provided should be used at least once. Do NOT skip images.
14. The "image_idx" field refers to the index in the images array provided. Match each image to the correct lesson content by reading its source slide title and context description carefully.
15. Avoid excessive use of bold/strong tags in text blocks. Use highlights sparingly for only the most important terms.

OUTPUT ONLY THE JSON ARRAY. No other text."""


def generate_slides_json(pdf_text, api_key, course_title=None, images_info=None, slide_text_notes=None):
    """Ask Claude to generate ONLY the slides JSON data."""
    client = Anthropic(api_key=api_key)

    # Don't truncate - send as much as possible (Sonnet has 200K context)
    if len(pdf_text) > 150000:
        pdf_text = pdf_text[:150000] + "\n\n[... Content truncated for context window ...]"

    title_instruction = f'Course Title: "{course_title}". Use this exact name for the course. ' if course_title else "Derive a clear, specific course title from the content (e.g. 'Investor Pitch Deck Masterclass', not 'Interactive Lesson'). "

    # Build image info section with rich context for intelligent placement
    images_section = ""
    if images_info:
        images_section = "\n\nAVAILABLE IMAGES — You MUST use ALL of these. Place each in the lesson slide matching its source topic:\n"
        for i, img in enumerate(images_info):
            images_section += f"  - image_idx {i}: {img['desc']}\n"
        images_section += "\nREMINDER: Match each image to the lesson content that covers the same topic. Write specific alt text describing what the image shows.\n"

    # Build per-slide text notes section
    notes_section = ""
    if slide_text_notes:
        notes_section = "\n\nUSER-ADDED NOTES FOR SPECIFIC SLIDES — incorporate these into the relevant lesson slides:\n"
        for slide_idx in sorted(slide_text_notes.keys()):
            notes_section += f"  - Slide {slide_idx + 1}: {slide_text_notes[slide_idx]}\n"
        notes_section += "\nThese are additional details the user wants included in the lesson. Weave them naturally into the content for the corresponding slide topics.\n"

    # Use streaming to avoid 10-minute timeout on long generations
    raw_chunks = []
    with client.messages.stream(
        model="claude-sonnet-4-20250514",
        max_tokens=32000,
        system=SLIDES_SYSTEM_PROMPT,
        messages=[{
            "role": "user",
            "content": f"""{title_instruction}Convert this PDF content into the JSON slides array.

IMPORTANT: Include ALL content from the PDF. Do NOT skip or summarize any sections. Every topic, concept, example, and detail must be covered. Create as many slides as needed.
{images_section}{notes_section}
PDF CONTENT:
{pdf_text}

Return ONLY the JSON array. No markdown, no explanation."""
        }],
    ) as stream:
        for text in stream.text_stream:
            raw_chunks.append(text)

    raw = "".join(raw_chunks).strip()

    # Strip markdown fences if present
    if raw.startswith("```"):
        lines = raw.split("\n")
        lines = lines[1:]  # remove opening fence
        if lines and lines[-1].strip().startswith("```"):
            lines = lines[:-1]
        raw = "\n".join(lines)

    return json.loads(raw)


def build_html(slides_data, course_title, elevenlabs_key="", elevenlabs_voice="", images=None):
    """Wrap the slides JSON in the complete, guaranteed-working HTML shell."""

    # Derive title if not provided — use first content slide's title
    if not course_title:
        for s in slides_data:
            if s.get("type") == "content" and s.get("t"):
                course_title = s["t"]
                break
        if not course_title:
            course_title = slides_data[0].get("t", "Lesson") if slides_data else "Lesson"

    # Build the welcome subtitle from first content slide
    welcome_sub = "Master the key concepts through interactive lessons, quizzes, and hands-on activities."
    for s in slides_data:
        if s.get("type") == "content" and s.get("s"):
            welcome_sub = s["s"]
            break

    slides_json = json.dumps(slides_data, ensure_ascii=False)

    # Build images lookup: index -> data_uri
    images_dict = {}
    if images:
        for i, img in enumerate(images):
            images_dict[i] = img["data_uri"]
    images_json = json.dumps(images_dict, ensure_ascii=False)

    return f'''<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1, maximum-scale=1, user-scalable=no">
<title>{course_title}</title>
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap" rel="stylesheet">
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
:root{{
  --b:#4E83FF;--b06:#4E83FF10;--b12:#4E83FF1f;--b25:#4E83FF40;
  --g:#16A34A;--g08:#16A34A14;--r:#DC2626;--r08:#DC262614;
  --y:#CA8A04;--y08:#CA8A0414;--gold:#F59E0B;
  --nv:#1A1F36;--co:#FF6B35;--co08:#FF6B3514;
  --s0:#FAFBFC;--s1:#F4F5F7;--s2:#E5E7EB;--s3:#D1D5DB;
  --c1:#111827;--c2:#4B5563;--c3:#9CA3AF;--c4:#C9CDD3;
  --rd:12px;
}}
body{{font-family:'Inter',system-ui,sans-serif;background:#fff;color:var(--c1);-webkit-font-smoothing:antialiased;font-weight:400;letter-spacing:-.01em;line-height:1.6;font-size:15px;overflow:hidden;height:100vh;overflow-wrap:break-word;word-wrap:break-word}}

@keyframes up{{from{{opacity:0;transform:translateY(24px)}}to{{opacity:1;transform:translateY(0)}}}}
@keyframes fadeIn{{from{{opacity:0}}to{{opacity:1}}}}
@keyframes slideR{{from{{opacity:0;transform:translateX(-20px)}}to{{opacity:1;transform:translateX(0)}}}}
@keyframes pop{{from{{opacity:0;transform:scale(.9)}}to{{opacity:1;transform:scale(1)}}}}
@keyframes slideDown{{from{{opacity:0;transform:translateY(-10px)}}to{{opacity:1;transform:translateY(0)}}}}
@keyframes xpFloat{{0%{{opacity:1;transform:translateY(0) scale(1)}}60%{{opacity:1;transform:translateY(-40px) scale(1.1)}}100%{{opacity:0;transform:translateY(-70px) scale(.9)}}}}
@keyframes xpPulse{{0%{{transform:scale(1)}}50%{{transform:scale(1.25)}}100%{{transform:scale(1)}}}}
@keyframes shake{{0%,100%{{transform:translateX(0)}}20%,60%{{transform:translateX(-4px)}}40%,80%{{transform:translateX(4px)}}}}
@keyframes glow{{0%,100%{{box-shadow:0 0 0 0 rgba(78,131,255,0)}}50%{{box-shadow:0 0 0 8px rgba(78,131,255,.15)}}}}
@keyframes correctBounce{{0%{{transform:scale(1)}}25%{{transform:scale(1.05)}}50%{{transform:scale(.97)}}75%{{transform:scale(1.02)}}100%{{transform:scale(1)}}}}
@keyframes wrongShake{{0%,100%{{transform:translateX(0)}}10%,50%,90%{{transform:translateX(-6px)}}30%,70%{{transform:translateX(6px)}}}}
@keyframes flashScreen{{0%{{opacity:.3}}100%{{opacity:0}}}}
@keyframes particle{{0%{{opacity:1;transform:translate(0,0) scale(1) rotate(0deg)}}100%{{opacity:0;transform:translate(var(--dx),var(--dy)) scale(0) rotate(var(--dr))}}}}
@keyframes starPop{{0%{{opacity:0;transform:scale(0) rotate(-30deg)}}50%{{opacity:1;transform:scale(1.4) rotate(5deg)}}100%{{opacity:0;transform:scale(.3) rotate(20deg) translateY(-20px)}}}}
@keyframes xpBoom{{0%{{transform:scale(1)}}30%{{transform:scale(1.5)}}60%{{transform:scale(.9)}}100%{{transform:scale(1)}}}}
@keyframes checkDraw{{to{{stroke-dashoffset:0}}}}

.an{{opacity:0}}.an.go{{animation:up .55s cubic-bezier(.16,1,.3,1) both}}
.an2{{opacity:0}}.an2.go{{animation:fadeIn .5s ease both}}
.an3{{opacity:0}}.an3.go{{animation:slideR .5s cubic-bezier(.16,1,.3,1) both}}
.an4{{opacity:0}}.an4.go{{animation:pop .45s cubic-bezier(.16,1,.3,1) both}}
.an5{{opacity:0}}.an5.go{{animation:slideDown .4s cubic-bezier(.16,1,.3,1) both}}

.app{{max-width:600px;margin:0 auto;height:100vh;display:flex;flex-direction:column;position:relative;overflow:hidden}}
.hd{{padding:16px 24px;display:flex;align-items:center;justify-content:space-between;background:rgba(255,255,255,.9);backdrop-filter:blur(16px);-webkit-backdrop-filter:blur(16px);z-index:50;border-bottom:1px solid var(--s2);flex-shrink:0}}
.hd-l{{display:flex;align-items:center;gap:12px}}
.ham{{background:none;border:none;cursor:pointer;padding:6px;border-radius:6px;transition:background .2s;display:flex;align-items:center}}
.ham:hover{{background:var(--s1)}}
.hd-cat{{font-size:13px;font-weight:500;color:var(--b);text-transform:uppercase;letter-spacing:1.5px}}
.hd-r{{display:flex;align-items:center;gap:14px}}
.hd-n{{font-size:12.5px;color:var(--c3)}}

.xp-badge{{display:flex;align-items:center;gap:5px;background:linear-gradient(135deg,#FEF3C7,#FDE68A);border:1px solid #FCD34D;border-radius:20px;padding:4px 12px 4px 8px;font-size:12.5px;color:#92400E;font-weight:500;position:relative;transition:all .3s}}
.xp-badge svg{{flex-shrink:0}}
.xp-toast{{position:absolute;top:-8px;right:-4px;font-size:12px;color:var(--gold);font-weight:500;pointer-events:none;animation:xpFloat 1.2s cubic-bezier(.16,1,.3,1) both;white-space:nowrap}}
.xp-pulse{{animation:xpPulse .4s cubic-bezier(.16,1,.3,1)}}

.bar{{height:2px;background:var(--s1);flex-shrink:0}}
.bar-f{{height:2px;background:var(--b);transition:width .6s cubic-bezier(.16,1,.3,1)}}

.ct{{flex:1;padding:28px 20px 24px;overflow-x:hidden;overflow-y:auto;-webkit-overflow-scrolling:touch}}
.ct.entering{{animation:slideEnter .4s cubic-bezier(.16,1,.3,1) both}}
@keyframes slideEnter{{from{{opacity:0;transform:translateX(30px)}}to{{opacity:1;transform:translateX(0)}}}}
@keyframes slideEnterBack{{from{{opacity:0;transform:translateX(-30px)}}to{{opacity:1;transform:translateX(0)}}}}
.ct.entering-back{{animation:slideEnterBack .4s cubic-bezier(.16,1,.3,1) both}}
.ct h1{{font-size:24px;font-weight:600;text-align:left;color:var(--c1);letter-spacing:-.3px;line-height:1.25;margin-bottom:6px}}
.ct .sub{{font-size:14px;color:var(--c2);line-height:1.6;margin-bottom:24px;text-align:left}}

.ft{{padding:16px 24px;background:rgba(255,255,255,.9);backdrop-filter:blur(16px);-webkit-backdrop-filter:blur(16px);z-index:50;border-top:1px solid var(--s2);display:flex;justify-content:space-between;align-items:center;flex-shrink:0}}
.bk{{background:none;border:none;font-size:14px;font-weight:500;color:var(--b);cursor:pointer;font-family:inherit;padding:8px 0;transition:opacity .2s}}
.bk:disabled{{color:var(--s3);cursor:default}}
.nx{{background:var(--nv);color:#fff;border:none;font-size:14px;font-weight:600;border-radius:12px;padding:12px 28px;cursor:pointer;font-family:inherit;transition:all .25s cubic-bezier(.16,1,.3,1)}}
.nx:hover{{transform:translateY(-1px);box-shadow:0 4px 12px rgba(0,0,0,.12)}}
.nx:disabled{{background:var(--s2);color:var(--c4);cursor:default;transform:none;box-shadow:none}}
.dots{{display:flex;gap:3px;align-items:center}}
.dt{{height:4px;border-radius:2px;transition:all .35s cubic-bezier(.16,1,.3,1);cursor:pointer}}
.dt.on{{width:16px;background:var(--b)}}
.dt.dn{{width:4px;background:var(--b25)}}
.dt.of{{width:4px;background:var(--s2)}}

.ov{{position:fixed;inset:0;background:rgba(0,0,0,.15);z-index:100;opacity:0;pointer-events:none;transition:opacity .25s}}
.ov.open{{opacity:1;pointer-events:auto}}
.dw{{position:fixed;top:0;left:0;bottom:0;width:264px;background:#fff;z-index:101;padding:28px 0;overflow-y:auto;transform:translateX(-100%);transition:transform .35s cubic-bezier(.16,1,.3,1)}}
.dw.open{{transform:translateX(0)}}
.dw-h{{padding:0 24px 20px;font-size:14px;font-weight:500;color:var(--c1)}}
.dw-c{{padding:12px 24px 4px;font-size:11px;font-weight:500;color:var(--c3);text-transform:uppercase;letter-spacing:1.5px}}
.dw-i{{display:flex;align-items:center;gap:8px;width:100%;padding:9px 24px 9px 28px;font-size:13.5px;color:var(--c2);background:transparent;border:none;text-align:left;cursor:pointer;font-family:inherit;transition:all .15s}}
.dw-i:hover{{background:var(--s0);color:var(--c1)}}
.dw-i.on{{color:var(--b);background:var(--b06)}}
.dw-i .dw-ico{{font-size:14px;width:20px;text-align:center}}

.crd{{background:#fff;border-radius:16px;border:1px solid var(--s2);padding:28px;max-width:480px;box-shadow:0 1px 3px rgba(0,0,0,.04);text-align:left;margin:0 auto}}
.ib{{border-radius:10px;padding:16px 20px;font-size:13px;line-height:1.65;text-align:left;max-width:440px;margin:14px auto}}
.ib.bl{{background:var(--b06);color:var(--c2)}}
.ib.gn{{background:var(--g08);color:var(--c2)}}
.ib.yw{{background:var(--y08);color:var(--c2)}}
.g2{{display:grid;grid-template-columns:1fr 1fr;gap:10px;max-width:100%}}
.pill{{border:none;border-radius:10px;padding:9px 16px;font-size:12.5px;font-weight:500;cursor:pointer;font-family:inherit;transition:all .2s}}
.pill.on{{background:var(--nv);color:#fff}}
.pill.of{{background:var(--s1);color:var(--c2)}}
.pill.of:hover{{background:var(--s2)}}

.qo{{background:var(--s0);border:1.5px solid var(--s2);border-radius:12px;padding:14px 18px;font-size:13.5px;color:var(--c1);text-align:left;cursor:pointer;font-family:inherit;width:100%;transition:all .25s cubic-bezier(.16,1,.3,1)}}
.qo:hover:not(:disabled){{border-color:var(--b);background:var(--b06)}}
.qo.ok{{background:var(--g08);border-color:var(--g);animation:correctBounce .5s cubic-bezier(.16,1,.3,1)}}
.qo.no{{background:var(--r08);border-color:var(--r);animation:wrongShake .5s ease}}
.xp-reward{{display:inline-flex;align-items:center;gap:4px;background:linear-gradient(135deg,#FEF3C7,#FDE68A);border-radius:20px;padding:3px 10px;font-size:12.5px;color:#92400E;font-weight:500;margin-left:6px}}

.cele-flash{{position:fixed;inset:0;pointer-events:none;z-index:200;animation:flashScreen .6s ease both}}
.cele-flash.green{{background:radial-gradient(circle at center,rgba(22,163,74,.25),transparent 70%)}}
.cele-flash.red{{background:radial-gradient(circle at center,rgba(220,38,38,.2),transparent 70%)}}
.particle{{position:fixed;pointer-events:none;z-index:201;border-radius:50%;animation:particle var(--dur) cubic-bezier(.16,1,.3,1) both}}
.particle.square{{border-radius:2px}}
.star-particle{{position:fixed;pointer-events:none;z-index:201;animation:starPop var(--dur) cubic-bezier(.16,1,.3,1) both}}

.check-circle{{display:inline-block;vertical-align:middle;margin-right:6px}}
.check-circle svg .check-path{{stroke-dasharray:20;stroke-dashoffset:20;animation:checkDraw .4s .2s ease both}}
.check-circle svg .circle-path{{stroke-dasharray:100;stroke-dashoffset:100;animation:checkDraw .5s ease both}}

.fas{{min-width:38px;height:38px;border-radius:50%;border:none;font-weight:500;font-size:13px;cursor:pointer;font-family:inherit;transition:all .3s cubic-bezier(.16,1,.3,1)}}
.fas.dn{{background:var(--g);color:#fff}}
.fas.nw{{background:var(--nv);color:#fff;animation:glow 2s ease infinite}}
.fas.wt{{background:var(--s1);color:var(--c3)}}

.tbb{{display:flex;gap:3px;background:var(--s1);border-radius:10px;padding:3px}}
.tbn{{flex:1;background:transparent;border:none;border-radius:8px;padding:9px 0;font-weight:500;font-size:12px;cursor:pointer;font-family:inherit;transition:all .25s cubic-bezier(.16,1,.3,1);color:var(--c3)}}
.tbn.on{{background:#fff;box-shadow:0 1px 4px rgba(0,0,0,.05);color:var(--c1)}}

.po{{background:#fff;color:var(--c2);border:1.5px solid var(--s2);border-radius:10px;padding:8px 14px;font-size:12px;cursor:pointer;font-family:inherit;transition:all .25s cubic-bezier(.16,1,.3,1)}}
.po:hover{{border-color:var(--b)}}
.po.on{{background:var(--nv);color:#fff;border-color:var(--nv)}}

.listen-badge{{display:flex;align-items:center;gap:5px;background:#F0FDFA;border:1px solid #99F6E4;border-radius:16px;padding:3px 10px 3px 6px;font-size:13px;color:#0D9488;font-weight:500;cursor:pointer;transition:all .2s}}
.listen-badge:hover{{background:#CCFBF1}}
.listen-badge .eq{{display:flex;align-items:flex-end;gap:1.5px;height:10px}}
.listen-badge .eq i{{width:2.5px;background:var(--b);border-radius:1px;animation:eqBar .8s ease infinite alternate}}
.listen-badge .eq i:nth-child(2){{animation-delay:.2s}}
.listen-badge .eq i:nth-child(3){{animation-delay:.4s}}
@keyframes eqBar{{from{{height:3px}}to{{height:10px}}}}
.listen-badge.off .eq i{{animation:none;height:3px}}
.listen-badge.off{{background:var(--s1);border-color:var(--s2);color:var(--c3)}}

@keyframes modalIn{{from{{opacity:0;transform:scale(.92) translateY(12px)}}to{{opacity:1;transform:scale(1) translateY(0)}}}}
@keyframes modalBgIn{{from{{opacity:0}}to{{opacity:1}}}}
.modal-bg{{position:fixed;inset:0;background:rgba(0,0,0,.25);backdrop-filter:blur(8px);-webkit-backdrop-filter:blur(8px);z-index:300;display:flex;align-items:center;justify-content:center;animation:modalBgIn .3s ease both}}
.modal{{background:#fff;border-radius:20px;padding:36px 28px;max-width:360px;width:90%;text-align:center;animation:modalIn .45s cubic-bezier(.16,1,.3,1) both;box-shadow:0 24px 60px rgba(0,0,0,.12)}}
.modal h2{{font-size:18px;font-weight:500;color:var(--c1);margin-bottom:6px}}
.modal p{{font-size:13px;color:var(--c2);line-height:1.6;margin-bottom:24px}}
.modal-btn{{width:100%;border:none;border-radius:12px;padding:14px;font-size:14px;font-weight:500;cursor:pointer;font-family:inherit;transition:all .25s cubic-bezier(.16,1,.3,1);margin-bottom:10px;display:flex;align-items:center;justify-content:center;gap:10px}}
.modal-btn:hover{{transform:translateY(-1px);box-shadow:0 4px 16px rgba(0,0,0,.1)}}
.modal-btn.primary{{background:var(--nv);color:#fff}}
.modal-btn.secondary{{background:var(--s1);color:var(--c2)}}
.modal-btn .btn-icon{{font-size:18px}}

.img-frame{{border:1px solid var(--s2);border-radius:var(--rd);overflow:hidden;margin:14px 0;background:var(--s0)}}
.img-frame img{{width:100%;display:block}}
.img-frame-label{{padding:8px 14px;font-size:12px;color:var(--c3);border-top:1px solid var(--s1)}}

@media(max-width:480px){{.g2{{grid-template-columns:1fr}}.ct{{padding:24px 18px 20px}}.hd,.ft{{padding:12px 18px}}}}
</style>
</head>
<body>
<div class="app" id="app"></div>
<script>
// ── DATA ──
const slidesData={slides_json};
const IMAGES={images_json};
const COURSE_TITLE=`{course_title}`;

// ── SVG CONSTANTS ──
const coinSvg=`<svg width="18" height="18" viewBox="0 0 512 512" fill="none" xmlns="http://www.w3.org/2000/svg"><circle cx="256" cy="256" r="256" fill="#fff"/><path d="M35.27 78.8c2.85-5.62 5.81-11 11.04-14.75 14.26-10.22 35.08-10.9 52.02-9.29 46.82 4.44 94 23.1 135.53 44.37 24.51 12.57 48.09 26.87 70.56 42.8 5.77 4.11 14.44 10.22 19.8 14.72 3.41 2.98 7.41 5.66 10.79 8.74.79.72 2.94 2.63 3.77 3.11 6.32-2.9 13.85-5.04 20.46-7.02 20.64-6.18 42.63-11.06 64.21-11.49 18.58-.37 42.62 2.35 51.44 21.3.42.91 1.03 2.01 1.58 2.84 1.16 1.63 1.72 3.15 2.72 4.7 2.76 4.32 5.09 8.81 7.67 13.22l14.16 24.5c7.32 12.75 12.81 19.38 9.89 35.2-2.42 13.12-8.08 22.62-14.88 33.76-3.01 4.63-8.02 11.83-11.83 15.75-1.05 1.59-1.97 2.57-3.19 4-9.51 11.12-20.16 21.58-30.91 31.49-2.32 2.1-11.93 10.76-14.06 11.69-1.55 1.46-3.23 2.87-4.93 4.16-10.66 8.1-21.8 17.08-33.01 24.4-2.14 1.31-13.78 9.3-15.06 9.68-12.55 7.9-27.52 16.94-40.76 23.58-1.34.97-14.13 7.46-16.12 8.2-11.28 5.28-22.46 10.62-34.09 15.1-3.17 1.23-6.46 2.43-9.57 3.77-2.5 1.29-14.06 5.12-17.05 6.02-17.53 5.31-29.48 8.78-47.93 11.97-3.89.8-16.46 2.39-20.02 2.08-10.47.66-18.56.24-28.86-1.82-21.21-4.24-25.55-15.37-35.34-32.4l-10.85-18.79c-3.36-5.91-7.09-12.75-10.67-18.43-.44-.98-2.43-4.3-3.13-5.4-3.06-4.87-4.53-10.34-4.59-16.11-.25-22.59 20.09-51.4 35.14-65.9-.63-.87-3.48-2.38-4.51-2.95-1.9-1.37-3.6-2.36-5.59-3.55-5.15-3.13-10.2-6.42-15.14-9.87-1.87-1.51-5.08-3.56-7.11-4.95-3.57-2.44-7.08-4.95-10.54-7.55-4.44-3.31-8.84-6.68-13.2-10.09-2.57-2.01-4.69-3.96-7.43-5.77-3.23-2.52-11.18-9.35-13.84-12.15-2.2-1.68-5.42-4.92-7.49-6.88-7.44-7.06-14.93-14.54-21.53-22.4-1.87-2.24-3.27-3.31-5.05-5.91-2.76-3.4-9.63-11.68-11.5-15.47C7.44 177.87-3.22 154.45 1.48 139.25c2.38-7.67 9.62-18.87 13.92-26.23 6.54-11.19 12.99-23.3 19.87-34.22z" fill="#FECD3E"/><path d="M35.27 78.8c2.85-5.62 5.81-11 11.04-14.75 14.26-10.22 35.08-10.9 52.02-9.29 46.82 4.44 94 23.1 135.53 44.37 24.51 12.57 48.09 26.87 70.56 42.8 5.77 4.11 14.44 10.22 19.8 14.72-.89.84-6.55 3.4-8.01 4.1-5.83 2.81-11.63 5.7-17.37 8.68-1.72.88-7.61 3.9-8.97 4.93l-.32-.05c-10.9 6.24-22.33 12.41-32.99 19.03-13.92 8.67-27.53 17.83-40.82 27.46l-.23.05c-8.65 6.57-17.38 12.94-25.67 19.97-1.25 1.06-3.03 2.68-4.37 3.51l.01.09c-11.52 9.16-25.39 22.71-35.5 33.24-2.78 2.95-5.5 5.96-8.14 9.04-1.43 1.65-5.68 6.97-7.15 7.9-1.89-1.36-3.59-2.36-5.59-3.55-5.15-3.13-10.2-6.42-15.14-9.87-1.87-1.51-5.08-3.56-7.11-4.95-3.57-2.44-7.08-4.95-10.54-7.55-4.44-3.31-8.84-6.68-13.2-10.09-2.57-2.01-4.69-3.96-7.43-5.77-3.23-2.52-11.18-9.35-13.84-12.15-2.2-1.68-5.42-4.92-7.49-6.88-7.44-7.06-14.93-14.54-21.53-22.4-1.87-2.24-3.27-3.31-5.05-5.91-2.76-3.4-9.63-11.68-11.5-15.47C7.44 177.87-3.22 154.45 1.48 139.25c2.38-7.67 9.62-18.87 13.92-26.23 6.54-11.19 12.99-23.3 19.87-34.22z" fill="#FECD3E"/><path d="M16.27 190.03C7.44 177.87-3.22 154.45 1.48 139.25c2.38-7.67 9.62-18.87 13.92-26.22 6.54-11.19 12.99-23.3 19.87-34.23.02.09.05.19.07.28.62 2.66-.23 5.47-.32 8.17-.2 6.4 1.41 13.03 3.54 19.02 16.66 46.69 80.87 96.59 122.47 123.15 5.06 3.23 10.16 6.39 15.32 9.46 2.79 1.66 6.5 3.7 9.14 5.46l.01.09c-11.52 9.16-25.39 22.71-35.5 33.24-2.78 2.95-5.5 5.96-8.14 9.04-1.43 1.65-5.68 6.97-7.15 7.9-1.89-1.36-3.59-2.36-5.59-3.55-5.15-3.13-10.2-6.42-15.14-9.87-1.87-1.51-5.08-3.56-7.11-4.95-3.57-2.44-7.08-4.95-10.54-7.54-4.44-3.32-8.84-6.68-13.2-10.09-2.57-2.01-4.69-3.96-7.43-5.77-3.23-2.52-11.18-9.36-13.84-12.15-2.2-1.68-5.42-4.92-7.49-6.88-7.44-7.06-14.93-14.54-21.53-22.4-1.87-2.24-3.27-3.31-5.05-5.9-2.76-3.4-9.63-11.68-11.5-15.48z" fill="#FEA02C"/><path d="M411.19 183.65c6.6-.63 16.5-.65 22.07 3.63 2.39 1.84 3.5 4.26 3.83 7.21.95 8.69-7.95 21.12-13.25 27.44C377.19 277.48 240.24 357.73 169.56 364.44c-6.71.41-16.38.66-21.82-4.11-2.09-1.83-3.36-4.41-3.54-7.18-.63-9.82 8.8-21.96 15.06-29.12C203.63 273.3 298.49 218.85 361.87 196.04c15.11-5.44 33.23-11.3 49.32-12.39z" fill="#FEA02C"/><path d="M365.18 81.86c1.6-.04 3.19-.11 4.79-.18-.12 5.9.19 12.78-.14 18.48 5.89-.2 12.35-.12 18.28-.16-.24 4.44-.06 10.24-.06 14.75-4.49-.02-14.6.24-18.45-.14.6 5.09.21 12.93.45 18.49-3.06-.14-6.62-.07-9.71-.08l-5.34.08c-.01-4.87-.25-14.01.18-18.57-5.18.71-12.45-.16-18.21.5-.07-5-.09-10.01-.05-15 5.66.15 12.94-.22 18.28.18-.43-4.49-.22-13.32-.24-18.17 3.4.03 6.79-.01 10.18-.11h.04zM34.1 298.96c4.86.14 10.09.03 14.98.02-.21 5.85.22 12.65-.19 18.27 5.69-.12 12.7.26 18.24-.26-.19 3.96-.19 11.19.13 15.04-5.31-.23-13.76.19-18.45-.22.2 1.91.45 17.16.13 18.4-4.45-.33-10.27-.13-14.82-.06-.08-1.3-.03-2.6-.01-3.9.1-4.85-.04-9.71.12-14.56-1.87.3-5.27.22-7.28.22-3.65-.02-7.29.03-10.94.15-.09-4.94-.03-10.13-.04-15.09 4.98.64 13.04.12 18.38.3-.35-3.48-.37-14.57-.28-18.3z" fill="#FEA02C"/></svg>`;
const logoSvg=`<svg width="28" height="28" viewBox="0 0 40 40" fill="none"><rect width="40" height="40" rx="10" fill="#4E83FF"/><path d="M10 28L20 12L30 28" stroke="#fff" stroke-width="2.5" stroke-linecap="round" stroke-linejoin="round"/><path d="M15 22L20 14L25 22" stroke="#fff" stroke-width="1.5" stroke-linecap="round" opacity=".5"/><circle cx="20" cy="10" r="2" fill="#fff"/></svg>`;
const animCheck=`<span class="check-circle"><svg width="20" height="20" viewBox="0 0 24 24" fill="none"><circle class="circle-path" cx="12" cy="12" r="10" stroke="#16A34A" stroke-width="2"/><path class="check-path" d="M8 12.5l2.5 3 5.5-6" stroke="#16A34A" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"/></svg></span>`;

// ── XP ──
let xp=0;
const COLORS=['#4E83FF','#003D5B','#FF6B35','#16A34A','#8B5CF6','#06B6D4','#EC4899'];
function addXP(n){{
  xp+=n;
  const badge=document.getElementById('xp-val');
  const wrap=document.getElementById('xp-wrap');
  if(badge){{badge.textContent=xp;wrap.classList.remove('xp-pulse');void wrap.offsetWidth;wrap.classList.add('xp-pulse')}}
  const toast=document.createElement('div');toast.className='xp-toast';toast.textContent=`+${{n}}`;
  if(wrap){{wrap.appendChild(toast);setTimeout(()=>toast.remove(),1300)}}
}}
function celebrate(originEl){{
  const flash=document.createElement('div');flash.className='cele-flash green';document.body.appendChild(flash);setTimeout(()=>flash.remove(),700);
  let cx=window.innerWidth/2,cy=window.innerHeight/2;
  if(originEl){{const r=originEl.getBoundingClientRect();cx=r.left+r.width/2;cy=r.top+r.height/2}}
  for(let i=0;i<24;i++){{const p=document.createElement('div');p.className='particle'+(Math.random()>.5?' square':'');const size=Math.random()*7+3;const angle=Math.random()*Math.PI*2;const dist=Math.random()*140+60;const dx=Math.cos(angle)*dist;const dy=Math.sin(angle)*dist-40;const dur=Math.random()*.4+.5;p.style.cssText=`left:${{cx}}px;top:${{cy}}px;width:${{size}}px;height:${{size}}px;background:${{COLORS[i%COLORS.length]}};--dx:${{dx}}px;--dy:${{dy}}px;--dr:${{Math.random()*400-200}}deg;--dur:${{dur}}s;`;document.body.appendChild(p);setTimeout(()=>p.remove(),dur*1000+50)}}
  for(let i=0;i<5;i++){{const s=document.createElement('div');s.className='star-particle';const ox=(Math.random()-.5)*120;const oy=(Math.random()-.5)*80-20;const dur=Math.random()*.3+.4;s.style.cssText=`left:${{cx+ox}}px;top:${{cy+oy}}px;--dur:${{dur}}s;`;s.innerHTML=`<svg width="16" height="16" viewBox="0 0 32 32" fill="${{COLORS[(i+3)%COLORS.length]}}"><ellipse cx="16" cy="16" rx="12" ry="10"/></svg>`;document.body.appendChild(s);setTimeout(()=>s.remove(),dur*1000+50)}}
}}
function wrongFlash(){{const flash=document.createElement('div');flash.className='cele-flash red';document.body.appendChild(flash);setTimeout(()=>flash.remove(),600)}}

// ── BUILD SLIDES ARRAY ──
function renderBlock(b){{
  if(!b) return '';
  if(typeof b==='string') return `<div class="an" style="font-size:13.5px;color:var(--c2);line-height:1.7;margin-bottom:12px">${{b}}</div>`;
  const k=b.kind||b.type||'';
  if(k==='text') return `<div class="an" style="font-size:13.5px;color:var(--c2);line-height:1.7;margin-bottom:12px">${{b.html||b.text||b.content||''}}</div>`;
  if(k==='bullets'){{
    const items=(b.items||[]).map(x=>`<li style="margin-bottom:6px">${{x}}</li>`).join('');
    return `<ul class="an" style="font-size:13.5px;color:var(--c2);line-height:1.7;padding-left:20px;margin-bottom:14px">${{items}}</ul>`;
  }}
  if(k==='icons'){{
    const items=(b.items||[]).map(x=>{{
      const label=x.label||x.text||x;
      const desc=x.desc||x.description||'';
      const icon=x.icon||x.emoji||'•';
      return `<div style="display:flex;align-items:flex-start;gap:10px;margin-bottom:10px"><span style="font-size:18px;flex-shrink:0">${{icon}}</span><div><div style="font-size:13.5px;color:var(--c2);line-height:1.6;font-weight:500">${{label}}</div>${{desc?`<div style="font-size:12.5px;color:var(--c3);line-height:1.5;margin-top:2px">${{desc}}</div>`:''}}</div></div>`;
    }}).join('');
    return `<div class="an" style="margin-bottom:14px">${{items}}</div>`;
  }}
  if(k==='steps'){{
    const items=(b.items||[]).map((x,i)=>{{
      const label=x.label||x.text||x;
      return `<div style="display:flex;align-items:flex-start;gap:12px;margin-bottom:12px"><div style="min-width:28px;height:28px;border-radius:50%;background:var(--b);color:#fff;display:flex;align-items:center;justify-content:center;font-size:12px;font-weight:600;flex-shrink:0">${{i+1}}</div><div style="font-size:13.5px;color:var(--c2);line-height:1.6;padding-top:4px">${{label}}</div></div>`;
    }}).join('');
    return `<div class="an" style="margin-bottom:14px">${{items}}</div>`;
  }}
  if(k==='tip'||k==='info'){{
    const cls=b.style==='green'?'gn':b.style==='yellow'?'yw':'bl';
    const icon=b.icon||b.emoji||(cls==='gn'?'\\u2705':cls==='yw'?'\\uD83D\\uDCA1':'\\u2139\\uFE0F');
    return `<div class="ib ${{cls}} an">${{icon}} &nbsp;${{b.text||b.content||''}}</div>`;
  }}
  if(k==='table'){{
    const hdr=(b.headers||[]).map(h=>`<th style="padding:10px 14px;text-align:left;font-size:12px;font-weight:500;color:var(--c3);text-transform:uppercase;letter-spacing:.5px;border-bottom:1px solid var(--s2)">${{h}}</th>`).join('');
    const rows=(b.rows||[]).map(row=>`<tr>${{row.map(cell=>`<td style="padding:10px 14px;font-size:13px;color:var(--c2);border-bottom:1px solid var(--s1)">${{cell}}</td>`).join('')}}</tr>`).join('');
    return `<div class="an" style="overflow-x:auto;margin-bottom:14px"><table style="width:100%;border-collapse:collapse;background:#fff;border:1px solid var(--s2);border-radius:10px;overflow:hidden"><thead><tr>${{hdr}}</tr></thead><tbody>${{rows}}</tbody></table></div>`;
  }}
  if(k==='code'){{
    return `<div class="an" style="margin-bottom:14px"><pre style="background:var(--nv);color:#e2e8f0;border-radius:10px;padding:18px;font-size:12.5px;line-height:1.6;overflow-x:auto;font-family:'Fira Code',monospace">${{b.code||b.text||b.content||''}}</pre></div>`;
  }}
  if(k==='compare'){{
    if(b.good!==undefined||b.bad!==undefined){{
      const gLabel=b.good_label||'Do This';const bLabel=b.bad_label||'Not This';
      return `<div class="g2 an" style="margin-bottom:14px"><div style="background:var(--g08);border-radius:10px;padding:14px 16px;font-size:13px;color:var(--c2);line-height:1.6"><strong>\\u2705 ${{gLabel}}</strong><br>${{b.good||''}}</div><div style="background:var(--r08);border-radius:10px;padding:14px 16px;font-size:13px;color:var(--c2);line-height:1.6"><strong>\\u274C ${{bLabel}}</strong><br>${{b.bad||''}}</div></div>`;
    }}
    const items=(b.items||[]).map(x=>{{
      const bg=x.color==='green'?'var(--g08)':x.color==='red'?'var(--r08)':'var(--b06)';
      const icon=x.icon||x.emoji||(x.color==='green'?'\\u2705':'\\u274C');
      return `<div style="background:${{bg}};border-radius:10px;padding:14px 16px;font-size:13px;color:var(--c2);line-height:1.6"><strong>${{icon}} ${{x.label||''}}</strong><br>${{x.text||x.content||''}}</div>`;
    }}).join('');
    return `<div class="g2 an" style="margin-bottom:14px">${{items}}</div>`;
  }}
  if(k==='image'){{
    const idx=b.image_idx;
    if(idx!==undefined && IMAGES[idx]){{
      const alt=b.alt||b.caption||'';
      return `<div class="img-frame an"><img src="${{IMAGES[idx]}}" alt="${{alt}}" loading="lazy">${{alt?`<div class="img-frame-label">${{alt}}</div>`:''}}</div>`;
    }}
    return '';
  }}
  if(k==='heading') return `<div class="an" style="font-size:16px;font-weight:600;color:var(--c1);margin:18px 0 8px">${{b.text||b.content||''}}</div>`;
  if(k==='divider') return `<hr class="an" style="border:none;border-top:1px solid var(--s2);margin:16px 0">`;
  // fallback: render any text-like content
  if(b.text||b.content) return `<div class="an" style="font-size:13.5px;color:var(--c2);line-height:1.7;margin-bottom:12px">${{b.text||b.content}}</div>`;
  return '';
}}

function buildContentSlide(d){{
  let html='<div style="max-width:100%">';
  const blocks=(d.body&&d.body.blocks)||d.body||[];
  if(Array.isArray(blocks)){{
    blocks.forEach(b=>{{ html+=renderBlock(b); }});
  }} else if(typeof blocks==='object'){{
    Object.values(blocks).forEach(b=>{{ if(Array.isArray(b))b.forEach(x=>{{html+=renderBlock(x)}}); }});
  }}
  html+='</div>';
  return html;
}}

const S=slidesData.map((d,idx)=>{{
  const obj={{cat:d.cat||'Lesson',t:d.t||'',s:d.s||'',narr:d.narration||d.narr||d.t+'. '+(d.s||'')}};
  const tp=d.type||'content';

  if(tp==='content'){{
    obj.r=function(){{return buildContentSlide(d)}};
  }}
  else if(tp==='quiz'){{
    const qid='q'+idx;
    const opts=d.options||(d.body&&d.body.options)||[];
    const ci=d.correct!==undefined?d.correct:(d.body&&d.body.correct!==undefined?d.body.correct:0);
    const q=d.question||(d.body&&d.body.question)||d.t;
    const exObj=d.explanations||(d.body&&d.body.explanations)||null;
    const ex=exObj?(typeof exObj==='string'?exObj:((exObj.correct||'')+(exObj.wrong?' '+exObj.wrong:''))):(d.explanation||(d.body&&d.body.explanation)||'');
    obj.r=function(){{return `<div id="${{qid}}" class="an"></div>`}};
    obj.i=function(){{QZ(qid,q,opts,ci,ex,true)}};
  }}
  else if(tp==='matching'){{
    const mid='m'+idx;
    const pairs=d.pairs||(d.body&&d.body.pairs)||[];
    obj.r=function(){{return `<div id="${{mid}}" class="an"></div>`}};
    obj.i=function(){{MATCH(mid,pairs)}};
  }}
  else if(tp==='ordering'){{
    const oid='o'+idx;
    const items=d.items||(d.body&&d.body.correct_order)||(d.body&&d.body.items)||[];
    obj.r=function(){{return `<div id="${{oid}}" class="an"></div>`}};
    obj.i=function(){{ORDER(oid,items)}};
  }}
  else if(tp==='prompt_builder'){{
    const pbid='pb'+idx;
    const rawParts=d.parts||(d.body&&d.body.parts)||[];
    const parts=(d.body&&d.body.chips)?[{{l:d.body.instructions||'Build your response',o:d.body.chips}}]:rawParts;
    obj.r=function(){{return `<div id="${{pbid}}" class="an"></div>`}};
    obj.i=function(){{PBUILD(pbid,parts)}};
  }}
  else if(tp==='milestone'){{
    const mEmoji=d.emoji||(d.body&&d.body.emoji)||'\\uD83C\\uDF89';
    const mMsg=d.s||(d.body&&d.body.message)||'Great progress! Keep going.';
    obj.r=function(){{
      return `<div style="text-align:center;padding:20px 0">
        <div class="an4" style="font-size:48px;margin-bottom:16px">${{mEmoji}}</div>
        <div class="an" style="font-size:20px;font-weight:600;color:var(--c1);margin-bottom:8px">${{d.t}}</div>
        <div class="an" style="font-size:14px;color:var(--c2);line-height:1.6;max-width:320px;margin:0 auto 20px">${{mMsg}}</div>
        <div class="an" style="display:inline-flex;align-items:center;gap:6px;background:linear-gradient(135deg,#FEF3C7,#FDE68A);border-radius:20px;padding:8px 20px;font-size:14px;color:#92400E;font-weight:500">${{coinSvg}} ${{xp}} XP earned</div>
      </div>`;
    }};
  }}
  else if(tp==='completion'){{
    const cEmoji=d.emoji||(d.body&&d.body.emoji)||'\\uD83C\\uDF93';
    const cMsg=d.s||(d.body&&d.body.message)||'You have completed the lesson. Well done!';
    obj.r=function(){{
      return `<div style="text-align:center;padding:20px 0">
        <div class="an4" style="font-size:56px;margin-bottom:16px">${{cEmoji}}</div>
        <div class="an" style="font-size:22px;font-weight:600;color:var(--c1);margin-bottom:8px">${{d.t||'Lesson Complete!'}}</div>
        <div class="an" style="font-size:14px;color:var(--c2);line-height:1.6;max-width:340px;margin:0 auto 24px">${{cMsg}}</div>
        <div class="an" style="display:inline-flex;align-items:center;gap:8px;background:linear-gradient(135deg,#FEF3C7,#FDE68A);border:2px solid #FCD34D;border-radius:24px;padding:12px 28px;font-size:18px;color:#92400E;font-weight:600">${{coinSvg}} ${{xp}} XP</div>
        <div class="an" style="margin-top:20px;font-size:12.5px;color:var(--c3)">Share your achievement or revisit any slide from the menu</div>
      </div>`;
    }};
  }}
  else{{
    obj.r=function(){{return buildContentSlide(d)}};
  }}
  return obj;
}});

// ── STATE ──
let cur=0,prevCur=0;
let listenMode=false;
let autoTimer=null;
let speaking=false;

// ── NAVIGATION ──
function go(i){{prevCur=cur;cur=Math.max(0,Math.min(S.length-1,i));stopAudio();R();if(listenMode)speakSlide()}}

// ── FOLLOW-ALONG STEPS ──
function FA(id,steps){{const el=document.getElementById(id);if(!el)return;let st=0;
  function r(){{const s=steps[st];el.innerHTML=`<div style="max-width:100%"><div style="display:flex;gap:8px;margin-bottom:20px;overflow-x:auto;padding-bottom:4px">${{steps.map((_,i)=>`<button class="fas ${{i<st?'dn':i===st?'nw':'wt'}}" onclick="window._f${{id}}(${{i}})">${{i<st?'\\u2713':i+1}}</button>`).join('')}}</div><div class="an4" style="background:var(--s0);border:1px solid var(--s2);border-radius:14px;padding:26px 22px"><div style="font-size:13px;font-weight:500;color:var(--b);text-transform:uppercase;letter-spacing:1.5px;margin-bottom:8px">Step ${{st+1}}</div><div style="font-size:15px;font-weight:500;color:var(--c1);margin-bottom:10px">${{s.t}}</div><div style="font-size:13px;color:var(--c2);line-height:1.65">${{s.d}}</div>${{s.p?`<div class="ib yw" style="margin-top:14px">\\uD83D\\uDCA1 &nbsp;${{s.p}}</div>`:''}}</div><div style="display:flex;justify-content:space-between;margin-top:16px"><button class="bk" onclick="window._f${{id}}(${{st-1}})" ${{st===0?'disabled':''}}>\\u2190 Previous</button><button class="nx" onclick="window._f${{id}}(${{st+1}})" ${{st===steps.length-1?'disabled':''}}>Next Step \\u2192</button></div></div>`;setTimeout(()=>{{const a=el.querySelector('.an4');if(a)a.classList.add('go')}},20)}}
  window['_f'+id]=function(i){{st=Math.max(0,Math.min(steps.length-1,i));r()}};r()}}

// ── PROMPT BUILDER ──
function PBUILD(id,parts){{const el=document.getElementById(id);if(!el)return;
  const pa=parts.map(p=>({{l:p.label||p.l||'Option',o:p.options||p.o||[]}}));
  let se=pa.map(()=>null);
  function r(){{const dn=se.every(s=>s!==null);el.innerHTML=`<div style="max-width:100%">${{pa.map((p,pi)=>`<div style="margin-bottom:18px"><div style="font-size:13px;font-weight:500;color:var(--b);text-transform:uppercase;letter-spacing:1.5px;margin-bottom:8px">${{p.l}}</div><div style="display:flex;gap:8px;flex-wrap:wrap">${{p.o.map((o,oi)=>`<button class="po${{se[pi]===oi?' on':''}}" onclick="window._pb${{id}}(${{pi}},${{oi}})">${{o}}</button>`).join('')}}</div></div>`).join('')}}${{dn?`<div class="an4 go" style="background:var(--s0);border:1px solid var(--s2);border-radius:10px;padding:18px"><div style="font-size:13px;font-weight:500;color:var(--c3);margin-bottom:6px;text-transform:uppercase;letter-spacing:1.5px">Your prompt</div><div style="font-size:13.5px;color:var(--c1);line-height:1.6;font-style:italic">"${{pa.map((p,i)=>p.o[se[i]]).join(', ')}}"</div></div>`:''}}</div>`}}
  window['_pb'+id]=function(p,o){{se[p]=o;r()}};r()}}

// ── MATCHING ──
function MATCH(id,pairs){{const el=document.getElementById(id);if(!el)return;
  const left=pairs.map((p,i)=>({{idx:i,text:p.left||p[0]||p.term||''}}));
  const right=pairs.map((p,i)=>({{idx:i,text:p.right||p[1]||p.def||''}}));
  // shuffle right side
  const shuffled=[...right].sort(()=>Math.random()-.5);
  let selL=null,matched={{}},wrongPair=null;
  function r(){{
    el.innerHTML=`<div class="crd"><div style="font-size:14px;font-weight:500;color:var(--c1);margin-bottom:16px">Match each item on the left with its pair on the right</div>
    <div style="display:grid;grid-template-columns:1fr 1fr;gap:10px">${{left.map((l,li)=>{{
      const isMatched=matched[l.idx]!==undefined;
      const isSelL=selL===li;
      const isWrongL=wrongPair&&wrongPair[0]===li;
      let cls='qo';if(isMatched)cls+=' ok';else if(isSelL)cls=' qo';else if(isWrongL)cls+=' no';
      let st=isSelL&&!isMatched?'border-color:var(--b);background:var(--b06)':'';
      return `<button class="${{cls}}" style="${{st}}" onclick="window._ml${{id}}(${{li}})" ${{isMatched?'disabled':''}}>${{isMatched?animCheck:''}}${{l.text}}</button>`;
    }}).join('')}}${{shuffled.map((r,ri)=>{{
      const isMatched=Object.values(matched).includes(ri);
      const isWrongR=wrongPair&&wrongPair[1]===ri;
      let cls='qo';if(isMatched)cls+=' ok';else if(isWrongR)cls+=' no';
      return `<button class="${{cls}}" onclick="window._mr${{id}}(${{ri}})" ${{isMatched?'disabled':''}}>${{r.text}}</button>`;
    }}).join('')}}</div>
    ${{Object.keys(matched).length===pairs.length?`<div class="an go" style="margin-top:14px;padding:14px;background:var(--g08);border-radius:10px;font-size:13px;color:var(--c1);text-align:center">${{animCheck}} All matched! <span class="xp-reward">${{coinSvg}} +20</span></div>`:''}}
    </div>`;
  }}
  window['_ml'+id]=function(li){{if(matched[left[li].idx]!==undefined)return;selL=li;r()}};
  window['_mr'+id]=function(ri){{
    if(selL===null)return;
    if(Object.values(matched).includes(ri))return;
    if(left[selL].idx===shuffled[ri].idx){{
      matched[left[selL].idx]=ri;selL=null;wrongPair=null;r();
      if(Object.keys(matched).length===pairs.length){{addXP(20);setTimeout(()=>celebrate(el),100)}}
    }}else{{
      wrongPair=[selL,ri];wrongFlash();r();setTimeout(()=>{{wrongPair=null;r()}},600);
    }}
    selL=null;
  }};
  r()}}

// ── ORDERING ──
function ORDER(id,items){{const el=document.getElementById(id);if(!el)return;
  const correct=items.map((x,i)=>i);
  let current=[...correct].sort(()=>Math.random()-.5);
  let selIdx=null,done=false;
  function r(){{
    el.innerHTML=`<div class="crd"><div style="font-size:14px;font-weight:500;color:var(--c1);margin-bottom:16px">Put these in the correct order</div>
    <div style="display:flex;flex-direction:column;gap:8px">${{current.map((ci,pos)=>{{
      const isSel=selIdx===pos;
      const item=typeof items[ci]==='string'?items[ci]:(items[ci].text||items[ci].label||items[ci]);
      let cls='qo';if(isSel)cls+=' ';
      let st=isSel?'border-color:var(--b);background:var(--b06)':'';
      if(done){{cls='qo ok';st='';}}
      return `<button class="${{cls}}" style="${{st}}" onclick="window._ord${{id}}(${{pos}})">${{done?animCheck:''}}${{pos+1}}. ${{item}}</button>`;
    }}).join('')}}</div>
    ${{done?`<div class="an go" style="margin-top:14px;padding:14px;background:var(--g08);border-radius:10px;font-size:13px;color:var(--c1);text-align:center">${{animCheck}} Correct order! <span class="xp-reward">${{coinSvg}} +20</span></div>`:
    `<button class="nx" style="margin-top:14px;width:100%" onclick="window._ordChk${{id}}()">Check Order</button>`}}
    </div>`;
  }}
  window['_ord'+id]=function(pos){{
    if(done)return;
    if(selIdx===null){{selIdx=pos;r()}}
    else{{const tmp=current[selIdx];current[selIdx]=current[pos];current[pos]=tmp;selIdx=null;r()}}
  }};
  window['_ordChk'+id]=function(){{
    const isCorrect=current.every((c,i)=>c===i);
    if(isCorrect){{done=true;addXP(20);r();setTimeout(()=>celebrate(el),100)}}
    else{{wrongFlash();const el2=document.getElementById(id);if(el2)el2.style.animation='wrongShake .5s ease';setTimeout(()=>{{if(el2)el2.style.animation=''}},600)}}
  }};
  r()}}

// ── QUIZ ──
function QZ(id,q,o,ci,ex,withXP){{const el=document.getElementById(id);if(!el)return;let sl=null;
  function r(){{const d=sl!==null;el.innerHTML=`<div class="crd" id="crd-${{id}}"><div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:16px"><div style="font-size:14px;color:var(--c1);line-height:1.55;flex:1">${{q}}</div>${{withXP?`<div style="font-size:12px;color:var(--gold);font-weight:500;margin-left:12px;white-space:nowrap">${{coinSvg}} 20</div>`:''}} </div><div style="display:flex;flex-direction:column;gap:8px">${{o.map((x,i)=>{{let c='qo';if(d&&i===ci)c+=' ok';if(d&&i===sl&&i!==ci)c+=' no';return`<button class="${{c}}" id="qo-${{id}}-${{i}}" onclick="window._q${{id}}(${{i}})" ${{d?'disabled':''}}>${{x}}</button>`}}).join('')}}</div>${{d?`<div class="an go" style="margin-top:14px;padding:14px;background:${{sl===ci?'var(--g08)':'var(--y08)'}};border-radius:10px;font-size:12.5px;color:var(--c1);line-height:1.6">${{sl===ci?`${{animCheck}} Correct! <span class="xp-reward">${{coinSvg}} +20</span><br>`:'\\u2717 Not quite. '}}${{ex}}</div>`:''}} </div>`}}
  window['_q'+id]=function(i){{if(sl===null){{sl=i;
    if(sl===ci){{if(withXP)addXP(20);r();setTimeout(()=>{{const btn=document.getElementById('qo-'+id+'-'+i);celebrate(btn)}},80)}}
    else{{wrongFlash();r()}}}}}};r()}}

// ── RENDER ──
function R(){{
  const s=S[cur],cats=[...new Set(S.map(x=>x.cat))],pct=((cur+1)/S.length)*100;
  let dots='';for(let i=0;i<S.length;i++)dots+=`<div class="dt ${{i===cur?'on':i<cur?'dn':'of'}}" onclick="go(${{i}})"></div>`;
  let nav='';cats.forEach(c=>{{nav+=`<div class="dw-c">${{c}}</div>`;S.filter(x=>x.cat===c).forEach(x=>{{const idx=S.indexOf(x);const ico=x.t.startsWith('Quick')?'\\u2726':'\\u2022';nav+=`<button class="dw-i${{idx===cur?' on':''}}" onclick="go(${{idx}});cN()"><span class="dw-ico">${{ico}}</span>${{x.t}}</button>`}})}});

  document.getElementById('app').innerHTML=`
    <div class="hd"><div class="hd-l"><button class="ham" onclick="oN()"><svg width="15" height="12" viewBox="0 0 15 12" fill="none"><path d="M1 1h13M1 6h9M1 11h13" stroke="var(--c1)" stroke-width="1.3" stroke-linecap="round"/></svg></button><span class="hd-cat">${{s.cat}}</span></div><div class="hd-r"><div id="listen-toggle" class="${{listenMode?'listen-badge':'listen-badge off'}}" onclick="toggleListen()"><div class="eq"><i></i><i></i><i></i></div><span class="listen-text">${{listenMode?'Listening':'Listen'}}</span></div><div class="xp-badge" id="xp-wrap">${{coinSvg}}<span id="xp-val">${{xp}}</span></div><span class="hd-n">${{cur+1}}/${{S.length}}</span></div></div>
    <div class="bar"><div class="bar-f" style="width:${{pct}}%"></div></div>
    <div class="ov" id="ov" onclick="cN()"></div><div class="dw" id="dw"><div class="dw-h">Lessons</div>${{nav}}</div>
    <div class="ct ${{cur>=prevCur?'entering':'entering-back'}}" id="cn"><h1 class="an">${{s.t}}</h1>${{s.s?`<p class="sub an">${{s.s}}</p>`:'<div style="height:20px"></div>'}}\n${{s.r()}}</div>
    <div class="ft"><button class="bk" onclick="go(${{cur-1}})" ${{cur===0?'disabled':''}}>\\u2190 Back</button><div class="dots">${{dots}}</div><button class="nx" onclick="go(${{cur+1}})" ${{cur===S.length-1?'disabled':''}}>Next \\u2192</button></div>`;

  setTimeout(()=>{{document.querySelectorAll('.an,.an2,.an3,.an4,.an5').forEach((el,i)=>{{setTimeout(()=>el.classList.add('go'),i*70)}})}},30);
  if(s.i)s.i();
  const cn=document.getElementById('cn');if(cn)cn.scrollTop=0;
}}
function oN(){{document.getElementById('ov').classList.add('open');document.getElementById('dw').classList.add('open')}}
function cN(){{document.getElementById('ov').classList.remove('open');document.getElementById('dw').classList.remove('open')}}

// ── TTS (ElevenLabs + fallback) ──
const EL_API_KEY='{elevenlabs_key}';
const EL_VOICE_ID='{elevenlabs_voice}';
const EL_MODEL='eleven_multilingual_v2';
let currentAudio=null;
let audioCache={{}};
let elAvailable=null;
let elTesting=false;
let cachedVoice=null;
const PRE_CACHE_AHEAD=3;

// Test ElevenLabs availability once at startup
async function testElevenLabs(){{
  if(elAvailable!==null)return elAvailable;
  if(elTesting)return false;
  if(!EL_API_KEY){{elAvailable=false;return false}}
  elTesting=true;
  try{{
    const resp=await fetch('https://api.elevenlabs.io/v1/user',{{
      method:'GET',headers:{{'xi-api-key':EL_API_KEY}}
    }});
    elAvailable=resp.ok;
  }}catch(e){{elAvailable=false}}
  elTesting=false;
  return elAvailable;
}}

async function fetchElevenLabsAudio(text,slideIdx){{
  if(audioCache[slideIdx])return audioCache[slideIdx];
  if(elAvailable===false)return null;
  try{{
    const resp=await fetch(`https://api.elevenlabs.io/v1/text-to-speech/${{EL_VOICE_ID}}/stream`,{{
      method:'POST',
      headers:{{'Content-Type':'application/json','xi-api-key':EL_API_KEY,'Accept':'audio/mpeg'}},
      body:JSON.stringify({{
        text:text,
        model_id:EL_MODEL,
        voice_settings:{{stability:0.5,similarity_boost:0.75,use_speaker_boost:true}}
      }})
    }});
    if(!resp.ok){{elAvailable=false;throw new Error('API '+resp.status)}}
    const blob=await resp.blob();
    const url=URL.createObjectURL(blob);
    audioCache[slideIdx]=url;
    return url;
  }}catch(e){{
    console.warn('Voice API unavailable:',e.message);
    elAvailable=false;
    return null;
  }}
}}

// Pre-cache multiple upcoming slides
function preCacheAhead(fromIdx){{
  for(let i=1;i<=PRE_CACHE_AHEAD;i++){{
    const idx=fromIdx+i;
    if(idx<S.length&&!audioCache[idx]){{
      const ns=S[idx];
      const text=ns.narr||ns.t+'. '+(ns.s||'');
      fetchElevenLabsAudio(text,idx).catch(()=>{{}});
    }}
  }}
}}

function stopAudio(){{
  if(currentAudio){{currentAudio.pause();currentAudio.currentTime=0;currentAudio=null}}
  if('speechSynthesis' in window)speechSynthesis.cancel();
  speaking=false;
  if(autoTimer){{clearTimeout(autoTimer);autoTimer=null}}
}}

function getVoice(){{
  if(cachedVoice)return cachedVoice;
  const voices=speechSynthesis.getVoices();
  if(!voices.length)return null;
  const neural=voices.filter(v=>v.lang.startsWith('en')&&(/natural|enhanced|premium|neural|online/i.test(v.name)));
  if(neural.length){{const fem=neural.find(v=>/aria|jenny|samantha|zira|siri|ava|alloy/i.test(v.name));if(fem){{cachedVoice=fem;return fem}}cachedVoice=neural[0];return neural[0]}}
  const prefs=['Samantha','Ava','Aria','Jenny','Karen','Moira','Tessa','Fiona','Google UK English Female','Google US English'];
  for(const p of prefs){{const v=voices.find(v=>v.name.includes(p)&&v.lang.startsWith('en'));if(v){{cachedVoice=v;return v}}}}
  const good=voices.filter(v=>v.lang.startsWith('en')&&!(/compact|espeak/i.test(v.name)));
  if(good.length){{cachedVoice=good[0];return good[0]}}
  cachedVoice=voices.find(v=>v.lang.startsWith('en'))||voices[0];
  return cachedVoice;
}}

// Split long text into chunks to avoid Chrome 15s speech cutoff bug
function splitTextForTTS(text){{
  if(text.length<200)return[text];
  const sentences=text.match(/[^.!?]+[.!?]+/g)||[text];
  const chunks=[];let chunk='';
  sentences.forEach(s=>{{
    if((chunk+s).length>180&&chunk){{chunks.push(chunk.trim());chunk=s}}
    else{{chunk+=s}}
  }});
  if(chunk.trim())chunks.push(chunk.trim());
  return chunks;
}}

function speakBrowserTTS(text,onEnd){{
  if(!('speechSynthesis' in window))return;
  const chunks=splitTextForTTS(text);
  let idx=0;
  function speakNext(){{
    if(idx>=chunks.length||!listenMode){{if(onEnd)onEnd();return}}
    const utter=new SpeechSynthesisUtterance(chunks[idx]);
    utter.voice=getVoice();utter.rate=0.95;utter.pitch=1.0;utter.volume=1;
    utter.onend=()=>{{idx++;speakNext()}};
    utter.onerror=(e)=>{{console.warn('Browser TTS error:',e);speaking=false}};
    speechSynthesis.speak(utter);
  }}
  // Delay first speak to avoid Chrome cancel/speak race condition
  setTimeout(speakNext,100);
}}

async function speakSlide(){{
  stopAudio();
  if(!listenMode)return;

  const myCur=cur;
  const s=S[myCur];
  let text=s.narr||s.t+'. '+(s.s||'');
  speaking=true;

  const badge=document.getElementById('listen-toggle');
  const setBadgeText=(t)=>{{if(badge){{const lt=badge.querySelector('.listen-text');if(lt)lt.textContent=t}}}};
  const stale=()=>!listenMode||cur!==myCur;

  const onSlideEnd=()=>{{
    speaking=false;currentAudio=null;
    if(stale())return;
    const isInteractive=s.t.startsWith('Quick Check')||s.t==='Build a Prompt';
    if(!isInteractive&&cur<S.length-1){{autoTimer=setTimeout(()=>{{go(cur+1)}},800)}}
  }};

  // Try ElevenLabs first (only if already confirmed available or not yet tested)
  if(EL_API_KEY&&elAvailable!==false){{
    // If cached, play instantly
    if(audioCache[myCur]){{
      setBadgeText('Listening');
      const audio=new Audio(audioCache[myCur]);
      currentAudio=audio;
      audio.onended=onSlideEnd;
      audio.onerror=()=>{{speaking=false;currentAudio=null}};
      try{{await audio.play()}}catch(e){{speaking=false;console.warn('Playback blocked:',e)}}
      preCacheAhead(myCur);
      return;
    }}
    // Not cached — fetch with loading indicator
    setBadgeText('Loading...');
    const audioUrl=await fetchElevenLabsAudio(text,myCur);
    if(stale()){{speaking=false;return}}

    if(audioUrl){{
      setBadgeText('Listening');
      const audio=new Audio(audioUrl);
      currentAudio=audio;
      audio.onended=onSlideEnd;
      audio.onerror=()=>{{speaking=false;currentAudio=null}};
      try{{await audio.play()}}catch(e){{speaking=false;console.warn('Playback blocked:',e)}}
      preCacheAhead(myCur);
      return;
    }}
  }}

  // Fallback: browser TTS
  if(stale()){{speaking=false;return}}
  setBadgeText('Listening');
  speakBrowserTTS(text,onSlideEnd);
}}

function toggleListen(){{
  listenMode=!listenMode;
  if(!listenMode){{stopAudio()}}
  else{{speakSlide()}}
  const badge=document.getElementById('listen-toggle');
  if(badge){{badge.className=listenMode?'listen-badge':'listen-badge off';badge.querySelector('.listen-text').textContent=listenMode?'Listening':'Listen'}}
}}

// Init: cache voice + test ElevenLabs at startup (non-blocking)
if('speechSynthesis' in window){{
  speechSynthesis.getVoices();
  speechSynthesis.onvoiceschanged=()=>{{speechSynthesis.getVoices();cachedVoice=null;getVoice()}}
  setTimeout(getVoice,100);
}}
if(EL_API_KEY)testElevenLabs().then(ok=>{{
  // Pre-cache first 3 slides if ElevenLabs is available
  if(ok)for(let i=0;i<Math.min(PRE_CACHE_AHEAD,S.length);i++){{
    const s=S[i];fetchElevenLabsAudio(s.narr||s.t+'. '+(s.s||''),i).catch(()=>{{}})
  }}
}})

// ── WELCOME MODAL ──
function showWelcome(){{
  const hasTTS='speechSynthesis' in window;
  const m=document.createElement('div');m.className='modal-bg';m.id='welcome-modal';
  m.innerHTML=`<div class="modal">
    <div style="margin-bottom:20px;font-size:40px">\\uD83D\\uDCDA</div>
    <h2>${{COURSE_TITLE}}</h2>
    <p>${{S[0]&&S[0].s?S[0].s:'Master key concepts through interactive lessons, quizzes, and activities.'}}</p>
    ${{hasTTS?`<button class="modal-btn primary" onclick="startListenMode()"><span class="btn-icon">\\uD83C\\uDFA7</span> Listen Along<span style="font-size:12.5px;color:rgba(255,255,255,.6);margin-left:4px">\\u00B7 auto-play</span></button>
    `:''}}
    <button class="modal-btn ${{hasTTS?'secondary':'primary'}}" onclick="startReadMode()"><span class="btn-icon">\\uD83D\\uDCD6</span> Read at My Pace</button>
    <div style="font-size:12px;color:var(--c3);margin-top:6px">${{S.length}} slides \\u00B7 Earn XP along the way</div>
  </div>`;
  document.body.appendChild(m);
}}
function startListenMode(){{listenMode=true;closeWelcome();speakSlide()}}
function startReadMode(){{listenMode=false;closeWelcome()}}
function closeWelcome(){{const m=document.getElementById('welcome-modal');if(m){{m.style.opacity='0';m.style.transition='opacity .25s';setTimeout(()=>m.remove(),260)}}}}

// ── KEYS ──
document.addEventListener('keydown',e=>{{if(e.key==='ArrowRight')go(cur+1);if(e.key==='ArrowLeft')go(cur-1)}});

R();
showWelcome();
</script>
</body>
</html>'''


def extract_pptx_slide_titles(filepath):
    """Extract slide titles from a PPTX file for the slide-by-slide image assignment UI."""
    prs = Presentation(filepath)
    slides = []
    for i, slide in enumerate(prs.slides):
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()
        if not title:
            # Try to get first text from any shape
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            title = text[:80]
                            break
                if title:
                    break
        slides.append({
            "index": i,
            "title": title or f"Slide {i + 1}",
            "slide_number": i + 1
        })
    return slides


def extract_pdf_page_titles(filepath):
    """Extract page titles/first-lines from a PDF for the slide-by-slide image assignment UI."""
    pages = []
    try:
        with pdfplumber.open(filepath) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                title = ""
                if text:
                    # Use first non-empty line as the title
                    for line in text.split("\n"):
                        line = line.strip()
                        if line:
                            title = line[:80]
                            break
                pages.append({
                    "index": i,
                    "title": title or f"Page {i + 1}",
                    "slide_number": i + 1
                })
    except Exception as e:
        print(f"PDF title extraction warning: {e}")
    return pages


def generate_lesson(pdf_text, api_key, course_title=None, elevenlabs_key="", elevenlabs_voice="", images=None, slide_text_notes=None):
    """Generate interactive HTML lesson from PDF text."""
    images_info = [{"page": img["page"], "desc": img["desc"]} for img in (images or [])]
    slides_data = generate_slides_json(pdf_text, api_key, course_title, images_info=images_info or None, slide_text_notes=slide_text_notes)
    title = course_title or "Interactive Lesson"
    return build_html(slides_data, title, elevenlabs_key=elevenlabs_key, elevenlabs_voice=elevenlabs_voice, images=images)


@app.route("/")
def index():
    return render_template("index.html")


@app.route("/extract", methods=["POST"])
def extract():
    """Extract slide titles from an uploaded PPTX/PDF without generating a lesson."""
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "No file selected"}), 400

    if not allowed_file(file.filename):
        return jsonify({"error": "Only PPTX and PDF files are allowed"}), 400

    filename = secure_filename(file.filename)
    ext = get_file_ext(filename)
    temp_path = os.path.join(app.config["UPLOAD_FOLDER"], f"temp_{uuid.uuid4().hex}_{filename}")
    file.save(temp_path)

    try:
        if ext in ("pptx", "ppt"):
            slides = extract_pptx_slide_titles(temp_path)
            file_type = "pptx"
        else:
            slides = extract_pdf_page_titles(temp_path)
            file_type = "pdf"

        return jsonify({
            "success": True,
            "file_type": file_type,
            "slides": slides,
            "total_slides": len(slides)
        })

    except Exception as e:
        return jsonify({"error": str(e)}), 500
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)


@app.route("/convert", methods=["POST"])
def convert():
    api_key = request.form.get("api_key", "").strip()
    if not api_key:
        return jsonify({"error": "Please provide your Anthropic API key"}), 400

    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "No file selected"}), 400

    if not allowed_file(file.filename):
        return jsonify({"error": "Only PPTX and PDF files are allowed"}), 400

    filename = secure_filename(file.filename)
    ext = get_file_ext(filename)
    temp_path = os.path.join(app.config["UPLOAD_FOLDER"], f"temp_{uuid.uuid4().hex}_{filename}")
    file.save(temp_path)

    try:
        # Extract text + images based on file type
        if ext in ("pptx", "ppt"):
            content_text = extract_pptx_text(temp_path)
            auto_images = extract_pptx_images(temp_path)
            slide_label = "Slide"
        else:
            content_text = extract_pdf_text(temp_path)
            auto_images = extract_pdf_images(temp_path)
            slide_label = "Page"

        if not content_text.strip():
            return jsonify({"error": "Could not extract text from this file. It may be empty or corrupted."}), 400

        # Process manually uploaded images (legacy bulk upload)
        manual_images = process_uploaded_images(request.files.getlist("images"))

        # Process slide-specific image assignments
        assigned_images = []
        image_assignments_json = request.form.get("image_assignments", "").strip()
        if image_assignments_json:
            try:
                assignments = json.loads(image_assignments_json)
                # assignments is a dict: {"0": "slide_image_0", "3": "slide_image_3", ...}
                # The keys are slide indices, values are the form field names
                for slide_idx_str, field_name in assignments.items():
                    slide_idx = int(slide_idx_str)
                    img_file = request.files.get(field_name)
                    if img_file and img_file.filename:
                        img_ext = get_file_ext(img_file.filename)
                        if img_ext in ("png", "jpg", "jpeg", "gif", "webp", "svg"):
                            blob = img_file.read()
                            content_type = img_file.content_type or f"image/{img_ext}"
                            b64 = base64.b64encode(blob).decode("utf-8")
                            assigned_images.append({
                                "page": slide_idx + 1,
                                "data_uri": f"data:{content_type};base64,{b64}",
                                "desc": f"User-assigned image for slide {slide_idx + 1}",
                                "source": "assigned"
                            })
            except (json.JSONDecodeError, ValueError) as e:
                print(f"Warning: Could not parse image_assignments: {e}")

        # Combine: assigned images first (highest priority), then manual uploads, then auto-extracted
        all_images = assigned_images + manual_images + auto_images

        # Process per-slide text notes
        slide_text_notes = {}
        text_notes_json = request.form.get("slide_text_notes", "").strip()
        if text_notes_json:
            try:
                raw_notes = json.loads(text_notes_json)
                # raw_notes is { "0": "text for slide 0", "3": "text for slide 3", ... }
                for slide_idx_str, text in raw_notes.items():
                    text = text.strip()
                    if text:
                        slide_text_notes[int(slide_idx_str)] = text
            except (json.JSONDecodeError, ValueError) as e:
                print(f"Warning: Could not parse slide_text_notes: {e}")

        course_title = request.form.get("title", "").strip()
        if not course_title:
            # Derive from filename: "Investor-Pitch-Deck-Lesson.pptx" -> "Investor Pitch Deck Lesson"
            raw_name = os.path.splitext(filename)[0]
            course_title = raw_name.replace("-", " ").replace("_", " ").strip()
            if not course_title:
                course_title = None
        elevenlabs_key = request.form.get("elevenlabs_key", "").strip()
        elevenlabs_voice = request.form.get("elevenlabs_voice", "").strip() or "21m00Tcm4TlvDq8ikWAM"
        html_content = generate_lesson(
            content_text, api_key, course_title,
            elevenlabs_key=elevenlabs_key, elevenlabs_voice=elevenlabs_voice,
            images=all_images, slide_text_notes=slide_text_notes
        )

        output_name = f"lesson_{uuid.uuid4().hex[:8]}.html"
        output_path = os.path.join(app.config["UPLOAD_FOLDER"], output_name)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(html_content)

        slides_count = content_text.count(f"--- {slide_label}")

        return jsonify({
            "success": True,
            "filename": output_name,
            "preview_url": f"/preview/{output_name}",
            "download_url": f"/download/{output_name}",
            "pages_extracted": slides_count,
            "chars_extracted": len(content_text),
            "images_extracted": len(auto_images),
            "images_uploaded": len(manual_images),
        })

    except json.JSONDecodeError as e:
        return jsonify({"error": f"Failed to parse AI response as JSON. Try again. Detail: {str(e)}"}), 500
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)


@app.route("/preview/<filename>")
def preview(filename):
    return send_from_directory(app.config["UPLOAD_FOLDER"], filename)


@app.route("/download/<filename>")
def download(filename):
    return send_from_directory(app.config["UPLOAD_FOLDER"], filename, as_attachment=True)


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5001))
    debug = not (os.environ.get("RAILWAY_ENVIRONMENT") or os.environ.get("VERCEL"))
    app.run(debug=debug, host="0.0.0.0", port=port)
